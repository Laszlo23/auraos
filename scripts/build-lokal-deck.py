#!/usr/bin/env python3
"""Build public/presentation-lokal.pptx from docs/deck-lokal screenshots."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "docs" / "deck-lokal"
OUT = ROOT / "public" / "presentation-lokal.pptx"

BG = RGBColor(0x0A, 0x0C, 0x12)
FG = RGBColor(0xF4, 0xF6, 0xFA)
MUTED = RGBColor(0x9A, 0xA3, 0xB2)
ACCENT = RGBColor(0x5B, 0xC4, 0xD4)
GOLD = RGBColor(0xD4, 0xA5, 0x4A)


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, size=28, bold=False, color=FG, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    return box


def title_slide(prs: Presentation, kicker: str, title: str, sub: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(0.4), kicker, size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.7), Inches(2.1), Inches(12), Inches(1.6), title, size=44, bold=True)
    add_textbox(slide, Inches(0.7), Inches(4.0), Inches(11), Inches(1.2), sub, size=18, color=MUTED)


def bullets_slide(prs: Presentation, kicker: str, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.35), kicker, size=13, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(12), Inches(1.0), title, size=34, bold=True)
    y = 2.3
    for b in bullets:
        add_textbox(slide, Inches(0.9), Inches(y), Inches(11.5), Inches(0.55), f"·  {b}", size=18, color=MUTED)
        y += 0.6


def shot_slide(prs: Presentation, kicker: str, title: str, caption: str, image: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(7), Inches(0.3), kicker, size=12, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.7), Inches(7), Inches(0.7), title, size=28, bold=True)
    add_textbox(slide, Inches(0.5), Inches(1.45), Inches(6.8), Inches(1.2), caption, size=15, color=MUTED)

    if image.exists():
        # Phone frame on the right
        left = Inches(8.15)
        top = Inches(0.55)
        height = Inches(6.5)
        # keep aspect ~390/844
        width = Emu(int(height * (390 / 844)))
        slide.shapes.add_picture(str(image), left, top, width=width, height=height)
    else:
        add_textbox(slide, Inches(8.2), Inches(3), Inches(4), Inches(1), f"Missing {image.name}", size=14, color=GOLD)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(
        prs,
        "AURA LOKAL",
        "Dein lokales Geschäft —\nonline, sichtbar, bewertet.",
        "Handy-App für Friseur, Beauty, Gastro & Immobilien.\nSocial · Kunden · Google-Bewertungen · Boost · Local Seat 99 €",
    )

    bullets_slide(
        prs,
        "01 · PROBLEM",
        "Heute ist lokal zu kompliziert.",
        [
            "Social-Kanäle, Bewertungen und Neukunden laufen getrennt.",
            "Agenturen sind teuer — Excel und WhatsApp sind chaotisch.",
            "Fake-Reviews sind riskant. Ehrliche Einladungen brauchen System.",
        ],
    )

    bullets_slide(
        prs,
        "02 · LÖSUNG",
        "Eine Super-App. Fünf Tabs.",
        [
            "Heute — nächster Schritt und Boost-Stand.",
            "Social — Kanäle verbinden, Posts freigeben.",
            "Kunden — Akquise & Outreach auf Deutsch.",
            "Bewertungen — Review Boost nur mit echten Kunden.",
            "Boost — Local Seat 99 € + klare Pakete.",
        ],
    )

    shot_slide(
        prs,
        "03 · LANDING",
        "/lokal",
        "Brand-first Einstieg für Service-Betriebe. CTA öffnet Auth mit funnel=local & lang=de.",
        SHOTS / "01-lokal-landing.png",
    )
    shot_slide(
        prs,
        "04 · HEUTE",
        "Home der App",
        "Ein Screen: nächster Schritt, Boost-Guthaben, Seat-Status und Shortcuts.",
        SHOTS / "05-heute.png",
    )
    shot_slide(
        prs,
        "05 · SOCIAL",
        "Kanäle & Posts",
        "OAuth verbinden, Entwürfe in Channels freigeben — ohne Agentur-Overhead.",
        SHOTS / "06-social.png",
    )
    shot_slide(
        prs,
        "06 · KUNDEN",
        "Neukunden & Outreach",
        "Akquise-Vorlagen für lokale Niches. Du bleibst Absender — Freigabe Pflicht.",
        SHOTS / "07-kunden.png",
    )
    shot_slide(
        prs,
        "07 · BEWERTUNGEN",
        "Google Review Boost",
        "Bis zu 999 Einladungen an echte Kunden. Track-Links, Klicks, founder-attestiert.",
        SHOTS / "08-bewertungen.png",
    )
    shot_slide(
        prs,
        "08 · BOOST",
        "Seat & Pakete",
        "99 € Local Seat (Barzahlung-Code oder Karte). Pakete: Sichtbarkeit · Bewertungen · Neukunden.",
        SHOTS / "09-boost.png",
    )
    shot_slide(
        prs,
        "09 · PUBLIC CARD",
        "/b/$slug",
        "Kundenkarte: Website-CTA + Google-Review-Button. Geteilt in einem Link.",
        SHOTS / "10-public-card.png",
    )

    bullets_slide(
        prs,
        "10 · COMPLIANCE",
        "Ehrlich bleiben.",
        [
            "Keine Fake-Reviews. Keine Bots, die als Kunden posten.",
            "Kein Scraping von Sterne-Zahlen als „Beweis“.",
            "Metriken = Einladungen · Klicks · founder-attestierte Reviews.",
            "Jeder Versand braucht Freigabe — wie bei Akquise.",
        ],
    )

    title_slide(
        prs,
        "NÄCHSTER SCHRITT",
        "Local Seat · 99 €",
        "Bar an der Theke (Code) oder Karte.\nStart: aibusiness.fun/lokal\nDeck: /presentation-lokal.pptx",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"saved {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
