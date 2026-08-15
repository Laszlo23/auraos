#!/usr/bin/env python3
"""Sales PowerPoint for Aura Lokal — German, for the person at the shop."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "public" / "Aura_Lokal_Verkauf.pptx"
MARK = ROOT / "public" / "brand" / "aura-mark.png"

BG = RGBColor(0x07, 0x09, 0x0E)
FG = RGBColor(0xF4, 0xF7, 0xFF)
MUTED = RGBColor(0x9A, 0xA3, 0xB2)
CYAN = RGBColor(0x4D, 0xE8, 0xF7)
GOLD = RGBColor(0xF4, 0xC0, 0x4A)
RED = RGBColor(0xC4, 0x45, 0x3A)

SLIDES = [
    {
        "kicker": "AURA LOKAL  ·  WIEN",
        "title": "Mehr echte Sterne.\nGäste, die wiederkommen.",
        "lines": [
            "Für Friseur, Beauty, Gastro, Handwerk, Immobilien.",
            "Kein Agentur-Theater. Keine Fake-Sterne.",
        ],
        "say": "Nicht das Betriebssystem verkaufen. Das Ergebnis: nach dem Besuch fragen — höflich, du gibst frei, Google bleibt ehrlich.",
    },
    {
        "kicker": "01  ·  DER ALLTAG",
        "title": "Gute Arbeit.\nStille Google-Seite.",
        "lines": [
            "Zufriedene Gäste gehen — und vergessen die Bewertung.",
            "Instagram, wenn’s einfällt. WhatsApp-Chaos. Eine Agentur, die teuer ist.",
            "Wer nicht fragt, bekommt keine Sterne. Wer falsch fragt, riskiert Google.",
        ],
        "say": "Frag: Wann hat zuletzt jemand von selbst bewertet? Die meisten sagen: selten. Genau da setzen wir an.",
    },
    {
        "kicker": "02  ·  WAS DER BETRIEB WILL",
        "title": "Drei Dinge.\nNicht fünfzig Features.",
        "lines": [
            "Mehr echte Google-Bewertungen — von Leuten, die wirklich da waren.",
            "Gäste, die wiederkommen, weil jemand nachfasst.",
            "Sichtbar bleiben, ohne eine Agentur zu füttern.",
        ],
        "say": "Warten bis der Mensch nickt. Dann erst der Preis. Wenn er über KI oder Token redet: zurück auf Sterne und Gäste.",
    },
    {
        "kicker": "03  ·  DAS ANGEBOT",
        "title": "Aura Reputation\n49 € im Monat.",
        "lines": [
            "Nach dem Besuch liegt eine höfliche Einladung bereit.",
            "Du tippst Freigeben. Fertig.",
            "Gäste können einchecken. Du bestätigst. Beziehungen bleiben lokal.",
        ],
        "say": "Ein Abo. Kein Paket-Salat am Tisch. Boost und Extra kommt später, wenn der Laden läuft.",
    },
    {
        "kicker": "04  ·  SO STARTET IHR",
        "title": "Drei Schritte.\nHeute.",
        "lines": [
            "1 · Kostenloser Check — Name, Stadt, Google-Link. Eine Minute.",
            "2 · Konto anlegen, Betrieb benennen.",
            "3 · Freischalten: 49 €/Monat mit Karte — oder Bar, Code an der Theke.",
        ],
        "say": "Am Tisch den Check machen. Handy des Betriebs, nicht deins. Screenshot vom Ergebnis. Dann zahlen oder Code.",
    },
    {
        "kicker": "05  ·  WAS DU BEKOMMST",
        "title": "Sterne. Gäste. Posts.",
        "lines": [
            "Sterne — Einladungen an echte Kunden. Klicks. Keine gekauften Texte.",
            "Gäste — QR oder Code. Du bestätigst den Besuch.",
            "Posts — Entwürfe warten. Du gibst frei, wenn’s passt.",
        ],
        "say": "Zeig Heute in der App, wenn du eingeloggt bist. Sonst diese drei Wörter. Nicht das ganze OS.",
    },
    {
        "kicker": "06  ·  WAS WIR NIE TUN",
        "title": "Keine Fake-Sterne.\nNiemals. Punkt.",
        "lines": [
            "Kein Geld für eine Google-Bewertung.",
            "Kein fertiger Text, den der Gast abschreiben soll.",
            "Keine Bots, keine gekauften Profile, keine fünf Sterne als Bedingung.",
        ],
        "say": "Wenn jemand nach gekauften Sternen fragt: oida, ned des. Google verbietet’s. Wien merkt sich Schmäh. Wir belohnen den Besuch — nicht die Sterne.",
    },
    {
        "kicker": "07  ·  PREIS",
        "title": "Klar. Klein.\nJeden Monat kündbar.",
        "lines": [
            "49 € im Monat mit Karte.",
            "129 € bar — ungefähr 3 Monate, Code an der Theke.",
            "Zuerst der Check. Der kostet nichts.",
        ],
        "say": "Nicht verhandeln unter 49. Bar ist der Code, nicht ein Rabatt. Wenn’s zu viel ist: Check dalassen, in einer Woche wiederkommen.",
    },
    {
        "kicker": "08  ·  WENN’S HAKT",
        "title": "Kurze Antworten.",
        "lines": [
            "Keine Zeit — du gibst nur frei. Den Rest legt Aura bereit.",
            "Ich zahl schon für Insta — das hier ist nach dem Besuch, nicht statt dem Feed.",
            "Bringt das was — wir zählen Einladungen und Klicks. Keine Fantasie-Sterne.",
            "Ist das erlaubt — ja, weil wir nicht für Reviews zahlen.",
        ],
        "say": "Ein Einwand, eine Antwort, dann zurück zum Check. Nicht diskutieren.",
    },
    {
        "kicker": "09  ·  HEUTE",
        "title": "Check machen.\nOder gleich starten.",
        "lines": [
            "aibusiness.fun/lokal/audit  —  eine Minute, am Tisch.",
            "aibusiness.fun/lokal  —  Konto, dann freischalten.",
            "aibusiness.fun/verkauf  —  diese Folien am Handy.",
        ],
        "say": "Schluss mit einer Handlung. Check öffnen oder Code schreiben. Nicht „ich schick dir was“ ohne Termin.",
    },
]


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_run(p, text: str, *, size=20, bold=False, color=FG, name="Calibri"):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    return run


def add_box(slide, left, top, width, height, text, *, size=28, bold=False, color=FG, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        add_run(p, line, size=size, bold=bold, color=color)
    return box


def add_notes(slide, text: str) -> None:
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = f"WAS DU SAGST\n{text}"


def austria_bar(slide, width) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(0.08))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    w = prs.slide_width

    for s in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        austria_bar(slide, w)
        if MARK.exists():
            slide.shapes.add_picture(str(MARK), Inches(0.55), Inches(0.28), Inches(0.38), Inches(0.38))
        add_box(slide, Inches(1.05), Inches(0.32), Inches(11), Inches(0.35), s["kicker"], size=13, bold=True, color=CYAN)
        add_box(slide, Inches(0.7), Inches(1.15), Inches(12), Inches(2.1), s["title"], size=40, bold=True, color=FG)
        y = 3.5
        for line in s["lines"]:
            add_box(slide, Inches(0.85), Inches(y), Inches(11.6), Inches(0.55), f"·   {line}", size=20, color=MUTED)
            y += 0.58
        add_notes(slide, s["say"])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"saved {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
